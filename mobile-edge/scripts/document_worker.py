"""Isolated document parser and generator used by the local PC gateway.

It receives only local paths and writes a JSON report.  The gateway owns
authentication, storage and download authorization; this worker owns format
parsing/generation and never fabricates an extraction result.
"""

from __future__ import annotations

import argparse
import hashlib
import json
import re
import sys
from pathlib import Path
from typing import Iterable

from docx import Document
from pptx import Presentation
from pypdf import PdfReader


MAX_CHUNK_CHARS = 1_400
OVERLAP_CHARS = 180


def chunks(text: str, locator: str) -> list[dict[str, str]]:
    clean = re.sub(r"\s+", " ", text).strip()
    if not clean:
        return []
    result: list[dict[str, str]] = []
    cursor = 0
    while cursor < len(clean):
        piece = clean[cursor : cursor + MAX_CHUNK_CHARS]
        result.append({"text": piece, "locator": locator, "sha256": hashlib.sha256(piece.encode("utf-8")).hexdigest()})
        if cursor + MAX_CHUNK_CHARS >= len(clean):
            break
        cursor += MAX_CHUNK_CHARS - OVERLAP_CHARS
    return result


def parse_document(source: Path) -> dict[str, object]:
    suffix = source.suffix.lower()
    if suffix == ".pdf":
        reader = PdfReader(str(source))
        values = []
        for number, page in enumerate(reader.pages, start=1):
            values.extend(chunks(page.extract_text() or "", f"第 {number} 页"))
        return {"parser_kind": "pypdf_text", "page_count": len(reader.pages), "chunks": values}
    if suffix == ".docx":
        document = Document(str(source))
        paragraphs = [paragraph.text for paragraph in document.paragraphs if paragraph.text.strip()]
        for table in document.tables:
            for row in table.rows:
                paragraphs.append(" | ".join(cell.text.strip() for cell in row.cells if cell.text.strip()))
        return {"parser_kind": "python_docx", "page_count": None, "chunks": chunks("\n".join(paragraphs), "正文与表格")}
    if suffix == ".pptx":
        presentation = Presentation(str(source))
        values: list[dict[str, str]] = []
        for number, slide in enumerate(presentation.slides, start=1):
            texts: list[str] = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    texts.append(shape.text)
                if getattr(shape, "has_table", False):
                    texts.extend(" | ".join(cell.text.strip() for cell in row.cells if cell.text.strip()) for row in shape.table.rows)
            values.extend(chunks("\n".join(texts), f"第 {number} 张幻灯片"))
        return {"parser_kind": "python_pptx", "page_count": len(presentation.slides), "chunks": values}
    if suffix in {".txt", ".md", ".csv", ".json"}:
        return {"parser_kind": "utf8_text", "page_count": None, "chunks": chunks(source.read_text(encoding="utf-8"), "正文")}
    raise ValueError("agent_resource_parser_unsupported")


def write_docx(path: Path, title: str, body: str, citations: Iterable[str]) -> None:
    document = Document()
    document.add_heading(title, 0)
    for paragraph in body.split("\n\n"):
        if paragraph.strip(): document.add_paragraph(paragraph.strip())
    cited = list(citations)
    if cited:
        document.add_heading("引用资料", level=1)
        for entry in cited: document.add_paragraph(entry, style="List Bullet")
    document.save(str(path))


def write_pptx(path: Path, title: str, body: str, citations: Iterable[str]) -> None:
    presentation = Presentation()
    cover = presentation.slides.add_slide(presentation.slide_layouts[0])
    cover.shapes.title.text = title
    cover.placeholders[1].text = "知行智学智能体生成"
    for index, paragraph in enumerate([item.strip() for item in body.split("\n\n") if item.strip()], start=1):
        slide = presentation.slides.add_slide(presentation.slide_layouts[1])
        slide.shapes.title.text = f"要点 {index}"
        slide.placeholders[1].text = paragraph
    cited = list(citations)
    if cited:
        slide = presentation.slides.add_slide(presentation.slide_layouts[1])
        slide.shapes.title.text = "引用资料"
        slide.placeholders[1].text = "\n".join(f"• {entry}" for entry in cited)
    presentation.save(str(path))


def write_pdf(path: Path, title: str, body: str, citations: Iterable[str]) -> None:
    from reportlab.lib.pagesizes import A4
    from reportlab.lib.styles import getSampleStyleSheet
    from reportlab.lib.units import cm
    from reportlab.pdfbase import pdfmetrics
    from reportlab.pdfbase.ttfonts import TTFont
    from reportlab.platypus import Paragraph, SimpleDocTemplate, Spacer

    font = Path("C:/Windows/Fonts/msyh.ttc")
    if not font.exists():
        raise ValueError("pdf_chinese_font_unavailable")
    pdfmetrics.registerFont(TTFont("ZhixingChinese", str(font), subfontIndex=0))
    styles = getSampleStyleSheet()
    styles["Title"].fontName = "ZhixingChinese"
    styles["BodyText"].fontName = "ZhixingChinese"
    story = [Paragraph(title, styles["Title"]), Spacer(1, 0.5 * cm)]
    for paragraph in body.split("\n\n"):
        if paragraph.strip(): story.extend([Paragraph(paragraph.strip().replace("\n", "<br/>"), styles["BodyText"]), Spacer(1, 0.25 * cm)])
    for entry in citations:
        story.append(Paragraph(f"引用：{entry}", styles["BodyText"]))
    SimpleDocTemplate(str(path), pagesize=A4, leftMargin=2 * cm, rightMargin=2 * cm).build(story)


def main() -> int:
    parser = argparse.ArgumentParser()
    sub = parser.add_subparsers(dest="command", required=True)
    parse = sub.add_parser("parse"); parse.add_argument("--input", type=Path, required=True); parse.add_argument("--output", type=Path, required=True)
    export = sub.add_parser("export"); export.add_argument("--format", choices=("docx", "pptx", "pdf"), required=True); export.add_argument("--input", type=Path, required=True); export.add_argument("--output", type=Path, required=True)
    args = parser.parse_args()
    try:
        if args.command == "parse":
            report = parse_document(args.input)
            report["state"] = "READY_FOR_AGENT" if report["chunks"] else "FAILED"
            report["error"] = None if report["chunks"] else "agent_resource_no_extractable_text"
            args.output.write_text(json.dumps(report, ensure_ascii=False), encoding="utf-8")
        else:
            source = json.loads(args.input.read_text(encoding="utf-8"))
            title, body, citations = source["title"], source["body"], source.get("citations", [])
            {"docx": write_docx, "pptx": write_pptx, "pdf": write_pdf}[args.format](args.output, title, body, citations)
            print(json.dumps({"state": "SUCCEEDED"}, ensure_ascii=False))
        return 0
    except Exception as error:
        print(json.dumps({"state": "FAILED", "error": str(error)[:240]}, ensure_ascii=False), file=sys.stderr)
        return 2


if __name__ == "__main__":
    raise SystemExit(main())
